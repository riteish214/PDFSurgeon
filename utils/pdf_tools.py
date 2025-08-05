"""
PDF manipulation utilities using various Python libraries
"""
import os
import tempfile
import logging
from PyPDF2 import PdfReader, PdfWriter
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import zipfile
import shutil
from io import BytesIO

logger = logging.getLogger(__name__)

class PDFTools:
    """
    Class containing all PDF manipulation methods
    """
    
    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
    
    def merge_pdfs(self, pdf_files):
        """
        Merge multiple PDF files into one
        
        Args:
            pdf_files (list): List of file paths to merge
            
        Returns:
            str: Path to the merged PDF file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'merged.pdf')
            writer = PdfWriter()
            
            for pdf_file in pdf_files:
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    writer.add_page(page)
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            logger.info(f"Successfully merged {len(pdf_files)} PDFs")
            return output_path
            
        except Exception as e:
            logger.error(f"Error merging PDFs: {str(e)}")
            raise Exception(f"Failed to merge PDFs: {str(e)}")
    
    def split_pdf(self, pdf_file):
        """
        Split PDF into individual pages
        
        Args:
            pdf_file (str): Path to the PDF file to split
            
        Returns:
            str: Path to directory containing split pages
        """
        try:
            reader = PdfReader(pdf_file)
            output_dir = os.path.join(self.temp_dir, 'split_pages')
            os.makedirs(output_dir, exist_ok=True)
            
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                
                output_path = os.path.join(output_dir, f'page_{i+1}.pdf')
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)
            
            logger.info(f"Successfully split PDF into {len(reader.pages)} pages")
            return output_dir
            
        except Exception as e:
            logger.error(f"Error splitting PDF: {str(e)}")
            raise Exception(f"Failed to split PDF: {str(e)}")
    
    def pdf_to_word(self, pdf_file):
        """
        Convert PDF to Word document
        
        Args:
            pdf_file (str): Path to the PDF file
            
        Returns:
            str: Path to the converted Word file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'converted.docx')
            
            # Use pdf2docx library for conversion
            cv = Converter(pdf_file)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            
            logger.info("Successfully converted PDF to Word")
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PDF to Word: {str(e)}")
            raise Exception(f"Failed to convert PDF to Word: {str(e)}")
    
    def pdf_to_ppt(self, pdf_file):
        """
        Convert PDF to PowerPoint presentation
        
        Args:
            pdf_file (str): Path to the PDF file
            
        Returns:
            str: Path to the converted PowerPoint file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'converted.pptx')
            
            # Create a new presentation
            prs = Presentation()
            
            # Read PDF and extract text/images for each page
            reader = PdfReader(pdf_file)
            
            for page_num, page in enumerate(reader.pages):
                # Add a slide for each page
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Extract text from the page
                text = page.extract_text()
                
                if text.strip():
                    # Add text box with extracted text
                    left = top = Inches(1)
                    width = Inches(8)
                    height = Inches(5.5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text[:1000]  # Limit text length
                
                # Add page number
                slide.shapes.title.text = f"Page {page_num + 1}"
            
            prs.save(output_path)
            
            logger.info("Successfully converted PDF to PowerPoint")
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PDF to PowerPoint: {str(e)}")
            raise Exception(f"Failed to convert PDF to PowerPoint: {str(e)}")
    
    def word_to_pdf(self, word_file):
        """
        Convert Word document to PDF
        Note: This is a simplified implementation
        
        Args:
            word_file (str): Path to the Word file
            
        Returns:
            str: Path to the converted PDF file
        """
        try:
            # This would typically require LibreOffice or similar
            # For now, we'll provide a placeholder implementation
            output_path = os.path.join(self.temp_dir, 'converted.pdf')
            
            # Placeholder: In production, you would use python-docx2pdf or LibreOffice
            # For this demo, we'll create a simple PDF with the filename
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            c.drawString(100, 750, f"Converted from: {os.path.basename(word_file)}")
            c.drawString(100, 730, "This is a demo conversion.")
            c.drawString(100, 710, "In production, this would contain the actual document content.")
            c.save()
            
            logger.info("Successfully converted Word to PDF (demo)")
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting Word to PDF: {str(e)}")
            raise Exception(f"Failed to convert Word to PDF: {str(e)}")
    
    def ppt_to_pdf(self, ppt_file):
        """
        Convert PowerPoint to PDF
        Note: This is a simplified implementation
        
        Args:
            ppt_file (str): Path to the PowerPoint file
            
        Returns:
            str: Path to the converted PDF file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'converted.pdf')
            
            # Placeholder implementation using reportlab
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            c.drawString(100, 750, f"Converted from: {os.path.basename(ppt_file)}")
            c.drawString(100, 730, "This is a demo conversion.")
            c.drawString(100, 710, "In production, this would contain the actual presentation content.")
            c.save()
            
            logger.info("Successfully converted PowerPoint to PDF (demo)")
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
            raise Exception(f"Failed to convert PowerPoint to PDF: {str(e)}")
    
    def compress_pdf(self, pdf_file):
        """
        Compress PDF file to reduce size
        
        Args:
            pdf_file (str): Path to the PDF file
            
        Returns:
            str: Path to the compressed PDF file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'compressed.pdf')
            
            reader = PdfReader(pdf_file)
            writer = PdfWriter()
            
            # Add all pages to writer
            for page in reader.pages:
                # Compress page content
                page.compress_content_streams()
                writer.add_page(page)
            
            # Write compressed PDF
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            # Calculate compression ratio
            original_size = os.path.getsize(pdf_file)
            compressed_size = os.path.getsize(output_path)
            ratio = (1 - compressed_size / original_size) * 100
            
            logger.info(f"Successfully compressed PDF by {ratio:.1f}%")
            return output_path
            
        except Exception as e:
            logger.error(f"Error compressing PDF: {str(e)}")
            raise Exception(f"Failed to compress PDF: {str(e)}")
    
    def rotate_pdf(self, pdf_file, rotation_angle, pages='all'):
        """
        Rotate pages in PDF
        
        Args:
            pdf_file (str): Path to the PDF file
            rotation_angle (int): Rotation angle (90, 180, 270)
            pages (str): Which pages to rotate ('all', 'odd', 'even')
            
        Returns:
            str: Path to the rotated PDF file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'rotated.pdf')
            
            reader = PdfReader(pdf_file)
            writer = PdfWriter()
            
            for i, page in enumerate(reader.pages):
                page_num = i + 1
                
                # Determine if this page should be rotated
                should_rotate = False
                if pages == 'all':
                    should_rotate = True
                elif pages == 'odd' and page_num % 2 == 1:
                    should_rotate = True
                elif pages == 'even' and page_num % 2 == 0:
                    should_rotate = True
                
                if should_rotate:
                    page.rotate(rotation_angle)
                
                writer.add_page(page)
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            logger.info(f"Successfully rotated PDF pages by {rotation_angle}°")
            return output_path
            
        except Exception as e:
            logger.error(f"Error rotating PDF: {str(e)}")
            raise Exception(f"Failed to rotate PDF: {str(e)}")
    
    def encrypt_pdf(self, pdf_file, password):
        """
        Add password protection to PDF
        
        Args:
            pdf_file (str): Path to the PDF file
            password (str): Password to protect the PDF
            
        Returns:
            str: Path to the encrypted PDF file
        """
        try:
            output_path = os.path.join(self.temp_dir, 'protected.pdf')
            
            reader = PdfReader(pdf_file)
            writer = PdfWriter()
            
            # Add all pages to writer
            for page in reader.pages:
                writer.add_page(page)
            
            # Encrypt with password
            writer.encrypt(password)
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            logger.info("Successfully encrypted PDF with password")
            return output_path
            
        except Exception as e:
            logger.error(f"Error encrypting PDF: {str(e)}")
            raise Exception(f"Failed to encrypt PDF: {str(e)}")
